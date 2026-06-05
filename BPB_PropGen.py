"""
Base Pair — Commercial Proposal Generator (Streamlit)
=====================================================

Turns an Otter.ai discovery-call transcript into a populated, non-binding
commercial proposal deck in Base Pair's standard format.

Pipeline:
    transcript  --(LLM extract)-->  structured parameters
    parameters  --(user review/edit in sidebar)-->  confirmed parameters
    parameters  --(LLM draft)-->  slide prose (Challenge / Strategy / Workflow)
    prose + manual milestones  --(review/edit)-->  final content
    final content --(template surgery)--> proposal_<customer>_<date>.pptx

Design notes:
  * Boilerplate slides 2-4 (aptamer tech, VennPlex SELEX, value prop) are NEVER
    touched -- they are company content that does not change project to project.
  * Only the project-specific slides are rewritten: slide 1 (title/subtitle/date),
    slide 5 (Challenge & Strategy), slide 6 (Workflow), slide 7 (Timeline,
    timeline step boxes, and the milestone/pricing table).
  * Shapes are addressed by stable shape_id (names are duplicated in the deck).
  * Pricing milestones are entered by the user -- never inferred from the call.
  * LLM provider is Anthropic (Claude). The entire LLM surface is one function,
    `call_llm()`; nothing else depends on the provider, so swapping is a one-spot
    change.

Run:  streamlit run basepair_proposal_app.py
"""

import os
import io
import json
import copy
import gzip
import base64
import zipfile
from datetime import datetime

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# --------------------------------------------------------------------------- #
# Constants
# --------------------------------------------------------------------------- #
APP_VERSION = "v21"
TEMPLATE_FILENAME = "proposal_template.pptx"
APP_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(APP_DIR, TEMPLATE_FILENAME)
LOGO_PATH = os.path.join(APP_DIR, "basepair_logo.png")
# Persisted API key lives in the user profile (NOT the project folder), so copying
# or sharing the app folder never leaks the key.
KEY_FILE = os.path.join(os.path.expanduser("~"), ".bpb_propgen_key")
PW_FILE = os.path.join(os.path.expanduser("~"), ".bpb_propgen_pw")

BP_GREEN_HEADER = RGBColor(0x54, 0x82, 0x35)   # table header
BP_GREEN_LIGHT  = RGBColor(0xE2, 0xEF, 0xDA)   # TOTAL row
WHITE           = RGBColor(0xFF, 0xFF, 0xFF)
DARK            = RGBColor(0x33, 0x33, 0x33)
TL_BOX_BLUE     = RGBColor(0x5B, 0x9B, 0xD5)   # timeline banner box fill

# Timeline banner geometry (inches), measured from the template's 4 boxes.
TL_BANNER_LEFT   = 0.519
TL_BANNER_RIGHT  = 9.328
TL_BANNER_TOP    = 1.593
TL_BOX_HEIGHT    = 0.816
TL_BOX_GAP       = 0.219      # horizontal gap between boxes

# Baked-in formatting defaults (the tweaks formerly applied by hand to every deck).
SUBTITLE_WIDTH_EMU   = 5379139            # slide 1: widen subtitle box to ~5.88"
TABLE_TOP_LIFT_EMU   = int(0.23 * 914400) # slide 7: raise the milestone table ~0.23"
TABLE_COL_FRACTIONS  = (0.544, 0.190, 0.266)  # Milestone / Pricing / Payment Terms
TABLE_HDR_ROW_IN     = 0.34               # header & TOTAL row height (inches)
TABLE_ROW_MIN_IN     = 0.40               # minimum milestone row height (inches)
TABLE_LINE_IN        = 0.19               # approx rendered height per wrapped line
TABLE_ROW_PAD_IN     = 0.10               # vertical padding added to each milestone row
TABLE_FOOT_GAP_EMU   = int(0.10 * 914400) # gap between table bottom and footnote

# Shape IDs of the variable shapes, keyed by 0-based slide index.
# (Boilerplate slides 2,3,4 -> indices 1,2,3 are intentionally absent.)
SLIDE_TITLE     = 0          # slide 1
SLIDE_CHALLENGE = 4          # slide 5
SLIDE_WORKFLOW  = 5          # slide 6
SLIDE_TIMELINE  = 6          # slide 7

ID_SUBTITLE   = 6            # slide 1: "Confidential Proposal to ... " textbox
ID_DATE       = 8            # slide 1: date textbox
ID_CHALLENGE  = 8            # slide 5: "Challenge: ..."
ID_STRATEGY   = 4            # slide 5: "Strategy: ..."
ID_WORKFLOW   = 3            # slide 6: bullet block
ID_TL_TITLE   = 2            # slide 7: title placeholder
ID_TL_LINE    = 13           # slide 7: "Timeline: ~6-9 weeks ..."
ID_TL_STEPS   = [31, 32, 33, 34]   # slide 7: the four timeline step boxes
ID_TABLE      = 20           # slide 7: milestone table
ID_TL_FOOT    = 4            # slide 7: footnote textbox

# Kd / affinity-method suggestions keyed by target type.
KD_METHOD_BLI = "Biolayer Interferometry (BLI)"
KD_METHOD_SMALL = "MicroScale Thermophoresis (MST) or Calorimetry"

# Suggested affinity (KD) method by target type:
#   Small molecule -> MST or Calorimetry; everything else -> BLI.
KD_METHOD_BY_TYPE = {
    "Protein":         KD_METHOD_BLI,
    "Peptide epitope": KD_METHOD_BLI,
    "Small molecule":  KD_METHOD_SMALL,
    "Whole cells":     KD_METHOD_BLI,
    "Viral particle":  KD_METHOD_BLI,
    "Other":           KD_METHOD_BLI,
}
TARGET_TYPES = list(KD_METHOD_BY_TYPE.keys())

# Options for the affinity-method dropdown.
KD_METHOD_OPTIONS = [
    KD_METHOD_BLI,
    KD_METHOD_SMALL,
    "MicroScale Thermophoresis (MST)",
    "Isothermal Titration Calorimetry (ITC)",
    "Surface Plasmon Resonance (SPR)",
    "Flow cytometry / fluorescence microscopy",
    "Other",
]


def parse_price(value):
    """Pull a number out of a free-text price like '$2,500', '14000', '1,400'."""
    import re
    if value is None:
        return None
    t = re.sub(r"[^0-9.\-]", "", str(value))
    if t in ("", "-", ".", "-."):
        return None
    try:
        return float(t)
    except ValueError:
        return None


def format_money(amount):
    """Format a number as USD with thousands separators and NO cents."""
    if amount is None:
        return "$0"
    return f"${int(round(amount)):,}"


def compute_total(milestones):
    """Sum the parseable milestone prices. Returns (numeric_total, n_parsed, n_unparsed)."""
    total, parsed, unparsed = 0.0, 0, 0
    for m in milestones:
        v = parse_price(m.get("price"))
        if v is None:
            if str(m.get("price", "")).strip():
                unparsed += 1
        else:
            total += v
            parsed += 1
    return total, parsed, unparsed


def parse_weeks(text):
    """Parse a weeks value that may be a single number or a range.
    Accepts '2', '2-3', '2–3', '~2 weeks', '1.5'. Returns (low, high) floats
    or None if no number is present."""
    if text is None:
        return None
    s = (str(text).lower().replace("weeks", "").replace("week", "")
         .replace("\u2013", "-").replace("\u2014", "-").replace("~", "").strip())
    nums = _re_weeks.findall(s)
    if not nums:
        return None
    if len(nums) == 1:
        v = float(nums[0]); return (v, v)
    return (float(nums[0]), float(nums[1]))


import re as _re_w
_re_weeks = _re_w.compile(r"\d+\.?\d*")


def _num(x):
    return str(int(x)) if float(x).is_integer() else f"{x:g}"


def normalize_weeks_cell(text):
    """Canonicalize a duration cell to bare number/range text: '2' or '2-3'.
    Strips words/symbols so the column stays purely numeric."""
    w = parse_weeks(text)
    if w is None:
        return ""
    lo, hi = w
    return _num(lo) if lo == hi else f"{_num(lo)}-{_num(hi)}"


def fmt_weeks(low, high):
    """Human label for the banner/timeline: '~2 weeks' or '~2-3 weeks'."""
    if low == high:
        unit = "week" if low == 1 else "weeks"
        return f"~{_num(low)} {unit}"
    return f"~{_num(low)}\u2013{_num(high)} weeks"


def compute_timeline(milestones):
    """Sum phase durations (weeks). Returns (low, high) or None if no durations."""
    lo = hi = 0.0
    any_ = False
    for m in milestones:
        w = parse_weeks(m.get("duration"))
        if w:
            lo += w[0]; hi += w[1]; any_ = True
    return (lo, hi) if any_ else None


def timeline_line(milestones):
    """Build the slide-7 'Timeline:' line from the summed phase weeks."""
    tl = compute_timeline(milestones)
    if tl is None:
        return "Timeline:       (enter phase durations in weeks below)"
    return f"Timeline:       {fmt_weeks(*tl)} from receipt of materials"


def distribute_ballpark(milestones, total, method="duration", round_to=100,
                        init_fraction=0.31):
    """Approximate each milestone price so they sum to `total`.
    The FIRST milestone (Project Initiation) is frontloaded at `init_fraction`
    of the total; the remainder is split across the other milestones, weighted
    by phase weeks (method='duration') or equally (method='even').
    Rounds to `round_to`; the largest of the remaining phases absorbs the
    rounding remainder so the grand total equals `total` exactly."""
    n = len(milestones)
    if n == 0 or total <= 0:
        return milestones
    if n == 1:
        milestones[0]["price"] = format_money(total)
        return milestones

    init = round((total * init_fraction) / round_to) * round_to
    remaining = total - init                      # exact: grand total stays = total
    rest = milestones[1:]

    if method == "duration":
        weights = []
        for m in rest:
            w = parse_weeks(m.get("duration"))
            weights.append((w[0] + w[1]) / 2 if w else 0.0)
        if sum(weights) == 0:
            weights = [1.0] * len(rest)
    else:
        weights = [1.0] * len(rest)
    s = sum(weights)
    raw = [remaining * w / s for w in weights]
    rounded = [round(r / round_to) * round_to for r in raw]
    diff = remaining - sum(rounded)
    idx = max(range(len(rest)), key=lambda i: raw[i])   # largest phase takes remainder
    rounded[idx] = max(0, rounded[idx] + diff)

    milestones[0]["price"] = format_money(init)
    for m, val in zip(rest, rounded):
        m["price"] = format_money(val)
    return milestones


def load_saved_key():
    """Resolve an Anthropic API key from, in order: env var, Streamlit secrets,
    or the local key file. Returns '' if none found."""
    env = os.getenv("ANTHROPIC_API_KEY")
    if env and env.strip():
        return env.strip()
    try:
        import streamlit as st
        val = st.secrets.get("ANTHROPIC_API_KEY", "")
        if val:
            return str(val).strip()
    except Exception:
        pass  # no secrets file configured
    try:
        if os.path.exists(KEY_FILE):
            with open(KEY_FILE, "r", encoding="utf-8") as fh:
                return fh.read().strip()
    except Exception:
        pass
    return ""


def save_key(key):
    """Persist the key to the user-profile key file. Returns True on success."""
    try:
        with open(KEY_FILE, "w", encoding="utf-8") as fh:
            fh.write(key.strip())
        try:
            os.chmod(KEY_FILE, 0o600)  # best-effort on POSIX; harmless on Windows
        except Exception:
            pass
        return True
    except Exception:
        return False


def forget_key():
    try:
        if os.path.exists(KEY_FILE):
            os.remove(KEY_FILE)
    except Exception:
        pass


def load_app_password():
    """Resolve the app gate password from, in order: Streamlit secrets
    (APP_PASSWORD), env var BPB_APP_PASSWORD, or a local password file.
    Returns '' if none configured. The password is NEVER stored in source so
    the repository can be public."""
    try:
        import streamlit as st
        val = st.secrets.get("APP_PASSWORD", "")
        if val:
            return str(val).strip()
    except Exception:
        pass
    env = os.getenv("BPB_APP_PASSWORD")
    if env and env.strip():
        return env.strip()
    try:
        if os.path.exists(PW_FILE):
            with open(PW_FILE, "r", encoding="utf-8") as fh:
                return fh.read().strip()
    except Exception:
        pass
    return ""


# --------------------------------------------------------------------------- #
# Low-level pptx helpers (no Streamlit dependency)
# --------------------------------------------------------------------------- #
def _shape_by_id(slide, shape_id):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    return None


def _clear_textframe(tf):
    """Remove all but one paragraph and empty it; reset its paragraph properties.

    The retained first paragraph's existing <a:pPr> is dropped so that bullet
    elements we add later are inserted into a clean element in valid schema order
    (a stale pPr containing defRPr/tabLst is what made PowerPoint flag the file)."""
    p_elements = tf._txBody.findall(qn("a:p"))
    for p in p_elements[1:]:
        tf._txBody.remove(p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        r._r.getparent().remove(r._r)
    existing_pPr = first._p.find(qn("a:pPr"))
    if existing_pPr is not None:
        first._p.remove(existing_pPr)
    return first


# Canonical child order inside <a:pPr> (CT_TextParagraphProperties). Inserting
# out of this order is what triggers PowerPoint's "repair" prompt.
_PPR_SEQ = ["a:lnSpc", "a:spcBef", "a:spcAft", "a:buClrTx", "a:buClr", "a:buSzTx",
            "a:buSzPct", "a:buSzPts", "a:buFontTx", "a:buFont", "a:buNone",
            "a:buAutoNum", "a:buChar", "a:tabLst", "a:defRPr", "a:extLst"]
_PPR_IDX = {qn(t): i for i, t in enumerate(_PPR_SEQ)}


def _insert_ordered(pPr, el):
    """Insert el into pPr at the position required by the OOXML schema."""
    idx = _PPR_IDX.get(el.tag, len(_PPR_SEQ))
    for child in pPr:
        if _PPR_IDX.get(child.tag, len(_PPR_SEQ)) > idx:
            child.addprevious(el)
            return
    pPr.append(el)


def _set_bullet(paragraph, char="\u2022", font="Arial", marL=228600, indent=-228600):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(marL))
    pPr.set("indent", str(indent))
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    _insert_ordered(pPr, pPr.makeelement(qn("a:buFont"), {"typeface": font}))
    _insert_ordered(pPr, pPr.makeelement(qn("a:buChar"), {"char": char}))


def _set_number(paragraph, marL=285750, indent=-285750):
    """Make the paragraph an auto-numbered list item (1. 2. 3. ...)."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(marL))
    pPr.set("indent", str(indent))
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    _insert_ordered(pPr, pPr.makeelement(qn("a:buFont"),
                                         {"typeface": "+mn-lt"}))
    _insert_ordered(pPr, pPr.makeelement(qn("a:buAutoNum"),
                                         {"type": "arabicPeriod"}))


def _no_bullet(paragraph, marL=0, indent=0):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(marL))
    pPr.set("indent", str(indent))
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    _insert_ordered(pPr, pPr.makeelement(qn("a:buNone"), {}))


def _add_run(paragraph, text, size_pt=14, bold=False, italic=False, color=None,
             font="Calibri"):
    run = paragraph.add_run()
    run.text = text
    f = run.font
    f.size = Pt(size_pt)
    f.bold = bold
    f.italic = italic
    f.name = font
    if color is not None:
        f.color.rgb = color
    return run


def _set_subscript(run):
    """Render a run as subscript (lowered baseline)."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("baseline", "-25000")   # -25% baseline = subscript


import re as _re
# The dissociation/affinity constant, written 'KD' or 'Kd' as a standalone token.
_KD_RE = _re.compile(r"\bK[Dd]\b")


def _emit(paragraph, text, size_pt=14, bold=False, italic=False, color=None,
          font="Calibri"):
    """Add text as runs, rendering the affinity constant 'KD'/'Kd' as K + subscript D."""
    kw = dict(size_pt=size_pt, bold=bold, italic=italic, color=color, font=font)
    if text is None:
        text = ""
    pos = 0
    matched = False
    for mo in _KD_RE.finditer(text):
        matched = True
        if mo.start() > pos:
            _add_run(paragraph, text[pos:mo.start()], **kw)
        _add_run(paragraph, "K", **kw)
        d_kw = dict(kw)
        d_kw["size_pt"] = max(7, round(size_pt * 0.72))
        _set_subscript(_add_run(paragraph, "D", **d_kw))   # always capital D
        pos = mo.end()
    if pos < len(text) or not matched:
        _add_run(paragraph, text[pos:], **kw)
    return paragraph


def _set_single_line(shape, text):
    """Replace a single-line textbox's text, preserving the first run's font."""
    tf = shape.text_frame
    first = _clear_textframe(tf)
    # try to inherit size from the original first run if it existed
    run = first.add_run()
    run.text = text


# --------------------------------------------------------------------------- #
# Slide writers
# --------------------------------------------------------------------------- #
def _write_title_slide(slide, p):
    sub = _shape_by_id(slide, ID_SUBTITLE)
    if sub is not None:
        sub.width = Emu(SUBTITLE_WIDTH_EMU)        # baked-in wider subtitle box
        tf = sub.text_frame
        para = _clear_textframe(tf)
        _no_bullet(para)
        _add_run(para, p["subtitle"], size_pt=28, bold=False, font="Calibri", color=WHITE)
        if p.get("phase_label"):
            blank = tf.add_paragraph()             # baked-in line return after title
            _no_bullet(blank)
            para2 = tf.add_paragraph()
            _no_bullet(para2)
            _add_run(para2, p["phase_label"], size_pt=18, italic=True,
                     font="Calibri", color=WHITE)
    date_sh = _shape_by_id(slide, ID_DATE)
    if date_sh is not None:
        para = _clear_textframe(date_sh.text_frame)
        _no_bullet(para)
        _add_run(para, p["proposal_date"], size_pt=24, font="Calibri", color=WHITE)


def _write_label_paragraph(shape, label, body, size_pt=14):
    """Bold lead-in label (e.g. 'Challenge:') followed by regular body text."""
    tf = shape.text_frame
    tf.word_wrap = True
    para = _clear_textframe(tf)
    _no_bullet(para)
    _add_run(para, f"{label}  ", size_pt=size_pt, bold=True, font="Calibri", color=DARK)
    _emit(para, body, size_pt=size_pt, bold=False, font="Calibri", color=DARK)


def _write_challenge_strategy(slide, p):
    ch = _shape_by_id(slide, ID_CHALLENGE)
    if ch is not None:
        _write_label_paragraph(ch, "Challenge:", p["challenge_text"], size_pt=14)
    stg = _shape_by_id(slide, ID_STRATEGY)
    if stg is not None:
        _write_label_paragraph(stg, "Strategy:", p["strategy_text"], size_pt=14)


def _write_workflow(slide, p):
    sh = _shape_by_id(slide, ID_WORKFLOW)
    if sh is None:
        return
    tf = sh.text_frame
    tf.word_wrap = True
    bullets = [b for b in p.get("workflow_bullets", []) if b.strip()]
    foot = p.get("workflow_footnote", "").strip()
    first = _clear_textframe(tf)
    size = 15
    for i, b in enumerate(bullets):
        para = first if i == 0 else tf.add_paragraph()
        _set_number(para)
        _emit(para, b, size_pt=size, font="Calibri", color=DARK)
    if foot:
        para = tf.add_paragraph()
        _no_bullet(para)
        _emit(para, foot, size_pt=11, italic=True, font="Calibri", color=DARK)


def _est_lines(text, col_width_emu, size_pt):
    """Estimate how many wrapped lines `text` needs in a column of the given
    width. Uses an average Calibri glyph advance of ~0.48 em."""
    col_in = int(col_width_emu) / 914400.0
    char_in = 0.48 * (size_pt / 72.0)          # average glyph width
    usable = max(col_in - 0.16, 0.5)           # subtract ~0.08" L/R cell insets
    per_line = max(int(usable / char_in), 1)
    n = len(str(text or ""))
    return max((n + per_line - 1) // per_line, 1)


def _rebuild_milestones_table(slide, p):
    """Delete the template table and add a fresh one sized to the milestones.

    Returns the table's bottom edge (Emu) so the caller can position the
    footnote beneath it regardless of how tall the table grew."""
    old = _shape_by_id(slide, ID_TABLE) or _find_table(slide)
    if old is None:
        return None
    left, top0, width = old.left, old.top, old.width
    top = Emu(int(top0) - TABLE_TOP_LIFT_EMU)        # baked-in: raise table ~0.23"
    # baked-in rebalanced columns (wider Payment Terms), as fractions of width
    col_w, acc = [], 0
    for k, frac in enumerate(TABLE_COL_FRACTIONS):
        w = int(width) - acc if k == len(TABLE_COL_FRACTIONS) - 1 else int(int(width) * frac)
        acc += w
        col_w.append(Emu(w))
    # remove the old table shape
    old._element.getparent().remove(old._element)

    milestones = p.get("milestones", [])
    rows = len(milestones) + 2          # header + milestones + TOTAL
    cols = 3

    # baked-in: content-driven row heights (replaces flat 0.42"/row that cramped
    # long descriptions). Header and TOTAL stay compact; milestone rows grow.
    row_h = [TABLE_HDR_ROW_IN]                       # header
    for m in milestones:
        lines = max(_est_lines(m.get("name", ""), col_w[0], 12),
                    _est_lines(m.get("terms", ""), col_w[2], 10))
        row_h.append(max(TABLE_ROW_MIN_IN, lines * TABLE_LINE_IN + TABLE_ROW_PAD_IN))
    row_h.append(TABLE_HDR_ROW_IN)                   # TOTAL

    height = Inches(sum(row_h))
    gf = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = gf.table
    gf.name = "BPB_MILESTONE_TABLE"     # stable tag for re-read / revise
    table.first_row = False             # we style the header ourselves
    table.horz_banding = False
    for j, w in enumerate(col_w):
        table.columns[j].width = w
    for i, h in enumerate(row_h):
        table.rows[i].height = Inches(h)

    headers = ["Milestone", "Pricing (USD)", "Payment Terms"]

    def style_cell(cell, text, bold=False, fill=None, color=DARK, size=12,
                   align=PP_ALIGN.LEFT):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill if fill else WHITE
        tf = cell.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = align
        _emit(para, text, size_pt=size, bold=bold, font="Calibri", color=color)

    # header
    for j, h in enumerate(headers):
        style_cell(table.cell(0, j), h, bold=True, fill=BP_GREEN_HEADER,
                   color=WHITE, size=12)
    # milestone rows
    for i, m in enumerate(milestones, start=1):
        price_raw = m.get("price", "")
        pv = parse_price(price_raw)
        price_disp = format_money(pv) if pv is not None else str(price_raw)
        style_cell(table.cell(i, 0), m.get("name", ""), size=12)
        style_cell(table.cell(i, 1), price_disp, size=12, align=PP_ALIGN.CENTER)
        style_cell(table.cell(i, 2), m.get("terms", ""), size=10)
    # total row
    last = rows - 1
    style_cell(table.cell(last, 0), "TOTAL", bold=True, fill=BP_GREEN_LIGHT, size=12)
    style_cell(table.cell(last, 1), p.get("total_price", ""), bold=True,
               fill=BP_GREEN_LIGHT, size=12, align=PP_ALIGN.CENTER)
    style_cell(table.cell(last, 2), "", fill=BP_GREEN_LIGHT)
    return Emu(int(top) + int(height))   # table bottom edge, for footnote placement


def _rebuild_timeline_boxes(slide, milestones):
    """Delete the template's 4 banner boxes and draw one box per milestone
    (4-6), evenly spaced across the banner. Each box shows the phase's short
    label (top) and duration (bottom)."""
    from pptx.enum.shapes import MSO_SHAPE
    # remove existing banner boxes: template ids 31-34 AND any tagged/edited boxes
    els = []
    for sid in ID_TL_STEPS:
        sh = _shape_by_id(slide, sid)
        if sh is not None:
            els.append(sh._element)
    for sh in _find_boxes(slide):
        els.append(sh._element)
    for el in {id(e): e for e in els}.values():
        el.getparent().remove(el)

    n = len(milestones)
    if n == 0:
        return
    span = TL_BANNER_RIGHT - TL_BANNER_LEFT
    # width so that n boxes + (n-1) gaps fill the span
    box_w = (span - (n - 1) * TL_BOX_GAP) / n
    label_sz = 12 if n <= 4 else (11 if n == 5 else 10)
    dur_sz = label_sz - 1
    for i, m in enumerate(milestones):
        left = TL_BANNER_LEFT + i * (box_w + TL_BOX_GAP)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(TL_BANNER_TOP),
            Inches(box_w), Inches(TL_BOX_HEIGHT))
        box.name = "BPB_TL_BOX"             # stable tag for re-read / revise
        box.fill.solid()
        box.fill.fore_color.rgb = TL_BOX_BLUE
        box.line.fill.background()
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        label = m.get("short_label", "").strip() or _derive_short_label(m.get("name", ""))
        _emit(para, label, size_pt=label_sz, bold=True, font="Calibri", color=WHITE)
        w = parse_weeks(m.get("duration", ""))
        if w:
            para2 = tf.add_paragraph()
            para2.alignment = PP_ALIGN.CENTER
            _add_run(para2, fmt_weeks(*w), size_pt=dur_sz, font="Calibri", color=WHITE)


def _derive_short_label(name, max_words=3):
    """Fallback banner label from a milestone name: first few significant words."""
    words = str(name).replace("&", "and").split()
    return " ".join(words[:max_words]) if words else ""


def _write_timeline_slide(slide, p):
    title = _shape_by_id(slide, ID_TL_TITLE)
    if title is not None and p.get("timeline_title"):
        para = _clear_textframe(title.text_frame)
        _add_run(para, p["timeline_title"], size_pt=28, bold=True, font="Calibri")
    line = _shape_by_id(slide, ID_TL_LINE)
    if line is not None:
        para = _clear_textframe(line.text_frame)
        _no_bullet(para)
        _add_run(para, p.get("timeline_estimate", ""), size_pt=18, bold=True,
                 font="Calibri", color=DARK)
    # banner boxes: one per milestone/phase
    _rebuild_timeline_boxes(slide, p.get("milestones", []))
    foot = _shape_by_id(slide, ID_TL_FOOT)
    if foot is not None and p.get("total_note"):
        para = _clear_textframe(foot.text_frame)
        _no_bullet(para)
        _add_run(para, p["total_note"], size_pt=10, italic=True, font="Calibri",
                 color=DARK)
    table_bottom = _rebuild_milestones_table(slide, p)
    if foot is not None and table_bottom is not None:
        foot.top = Emu(int(table_bottom) + TABLE_FOOT_GAP_EMU)  # follow table bottom


# --------------------------------------------------------------------------- #
# Public engine entry point
# --------------------------------------------------------------------------- #
def build_proposal_pptx(params, template_path=TEMPLATE_FILENAME):
    """Return the populated proposal deck as bytes. Streamlit-free.

    The full parameter set is embedded in the file's custom document properties
    so the deck can later be re-uploaded and revised (see revise_proposal_pptx)."""
    prs = Presentation(template_path)
    slides = list(prs.slides)
    _write_title_slide(slides[SLIDE_TITLE], params)
    _write_challenge_strategy(slides[SLIDE_CHALLENGE], params)
    _write_workflow(slides[SLIDE_WORKFLOW], params)
    _write_timeline_slide(slides[SLIDE_TIMELINE], params)
    bio = io.BytesIO()
    prs.save(bio)
    state = dict(params)
    state.setdefault("_revision", 0)
    return _embed_state_in_bytes(bio.getvalue(), state)


# --------------------------------------------------------------------------- #
# Revision engine: build the next version ONTO the user's edited deck so that
# manual formatting and minor text edits are preserved across revisions.
# --------------------------------------------------------------------------- #

# ---- embedded state (custom document properties survive a PowerPoint round-trip)
_CUSTOM_NS = "http://schemas.openxmlformats.org/officeDocument/2006/custom-properties"
_VT_NS = "http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"
_STATE_PROP = "BPBState"
_STATE_FMTID = "{D5CDD505-2E9C-101B-9397-08002B2CF9AE}"


def _encode_state(params):
    blob = json.dumps(params, ensure_ascii=False).encode("utf-8")
    return base64.b64encode(gzip.compress(blob)).decode("ascii")


def _decode_state(text):
    try:
        return json.loads(gzip.decompress(base64.b64decode(text)).decode("utf-8"))
    except Exception:
        return None


def _embed_state_in_bytes(pptx_bytes, params):
    """Write params into docProps/custom.xml (property 'BPBState') and return new bytes."""
    zin = zipfile.ZipFile(io.BytesIO(pptx_bytes), "r")
    items = {n: zin.read(n) for n in zin.namelist()}
    zin.close()

    C = "{%s}" % _CUSTOM_NS
    V = "{%s}" % _VT_NS
    if "docProps/custom.xml" in items:
        root = etree.fromstring(items["docProps/custom.xml"])
        for pr in list(root.findall(C + "property")):
            if pr.get("name") == _STATE_PROP:
                root.remove(pr)
        pids = [int(pr.get("pid", "1")) for pr in root.findall(C + "property")]
        next_pid = (max(pids) if pids else 1) + 1
    else:
        root = etree.Element(C + "Properties", nsmap={None: _CUSTOM_NS, "vt": _VT_NS})
        next_pid = 2
    prop = etree.SubElement(root, C + "property",
                            {"fmtid": _STATE_FMTID, "pid": str(next_pid), "name": _STATE_PROP})
    etree.SubElement(prop, V + "lpwstr").text = _encode_state(params)
    items["docProps/custom.xml"] = etree.tostring(
        root, xml_declaration=True, encoding="UTF-8", standalone=True)

    ct = items.get("[Content_Types].xml", b"")
    if b"docProps/custom.xml" not in ct:
        ins = (b'<Override PartName="/docProps/custom.xml" ContentType='
               b'"application/vnd.openxmlformats-officedocument.custom-properties+xml"/>')
        items["[Content_Types].xml"] = ct.replace(b"</Types>", ins + b"</Types>")

    rels = items.get("_rels/.rels", b"")
    if b"custom-properties" not in rels:
        rel = (b'<Relationship Id="rIdBPBState" Type='
               b'"http://schemas.openxmlformats.org/officeDocument/2006/relationships/'
               b'custom-properties" Target="docProps/custom.xml"/>')
        items["_rels/.rels"] = rels.replace(b"</Relationships>", rel + b"</Relationships>")

    out = io.BytesIO()
    zo = zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED)
    for n, data in items.items():
        zo.writestr(n, data)
    zo.close()
    return out.getvalue()


def _read_state_from_bytes(pptx_bytes):
    try:
        zin = zipfile.ZipFile(io.BytesIO(pptx_bytes), "r")
        if "docProps/custom.xml" not in zin.namelist():
            zin.close(); return None
        xml = zin.read("docProps/custom.xml"); zin.close()
        root = etree.fromstring(xml)
        C = "{%s}" % _CUSTOM_NS; V = "{%s}" % _VT_NS
        for pr in root.findall(C + "property"):
            if pr.get("name") == _STATE_PROP:
                node = pr.find(V + "lpwstr")
                if node is not None and node.text:
                    return _decode_state(node.text)
    except Exception:
        return None
    return None


# ---- shape locators that work on edited decks (template ids no longer apply) ----
def _find_table(slide):
    for sh in slide.shapes:
        if sh.has_table:
            return sh
    return None


def _find_boxes(slide):
    tagged = [sh for sh in slide.shapes if (sh.name or "") == "BPB_TL_BOX"]
    if tagged:
        return tagged
    out = []
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                if sh.text_frame.text.strip():
                    out.append(sh)
            except Exception:
                pass
    return out


# ---- formatting samplers / format-preserving writers ----
def _sample_fmt(shape_or_cell):
    """Read size/color/font/bold/italic from the first run of a shape or cell."""
    out = dict(size_pt=14, color=DARK, font="Calibri", bold=False, italic=False)
    try:
        for p in shape_or_cell.text_frame.paragraphs:
            if p.runs:
                f = p.runs[0].font
                if f.size is not None:
                    out["size_pt"] = f.size.pt
                out["font"] = f.name or "Calibri"
                out["bold"] = bool(f.bold)
                out["italic"] = bool(f.italic)
                try:
                    if f.color is not None and f.color.type is not None:
                        out["color"] = f.color.rgb
                except Exception:
                    pass
                break
    except Exception:
        pass
    return out


def _set_text_preserve(shape, text):
    if shape is None:
        return
    fmt = _sample_fmt(shape)
    align = shape.text_frame.paragraphs[0].alignment
    para = _clear_textframe(shape.text_frame)
    _no_bullet(para)
    if align is not None:
        para.alignment = align
    _emit(para, text, size_pt=fmt["size_pt"], bold=fmt["bold"], italic=fmt["italic"],
          color=fmt["color"], font=fmt["font"])


def _set_label_body_preserve(shape, label, body):
    if shape is None:
        return
    fmt = _sample_fmt(shape)            # font size = whatever the user set (e.g. 13pt)
    tf = shape.text_frame
    tf.word_wrap = True
    para = _clear_textframe(tf)
    _no_bullet(para)
    _add_run(para, f"{label}  ", size_pt=fmt["size_pt"], bold=True,
             font=fmt["font"], color=fmt["color"])
    _emit(para, body, size_pt=fmt["size_pt"], bold=False, font=fmt["font"], color=fmt["color"])


def _rewrite_subtitle_preserve(shape, params):
    if shape is None:
        return
    tf = shape.text_frame
    sub_fmt = _sample_fmt(shape)
    phase_size = 18
    for p in tf.paragraphs:
        for r in p.runs:
            if r.font.italic:
                if r.font.size is not None:
                    phase_size = r.font.size.pt
                break
    para = _clear_textframe(tf)
    _no_bullet(para)
    _add_run(para, params.get("subtitle", ""), size_pt=sub_fmt["size_pt"],
             color=sub_fmt["color"], font=sub_fmt["font"])
    if params.get("phase_label"):
        blank = tf.add_paragraph(); _no_bullet(blank)
        p2 = tf.add_paragraph(); _no_bullet(p2)
        _add_run(p2, params["phase_label"], size_pt=phase_size, italic=True,
                 color=sub_fmt["color"], font=sub_fmt["font"])


def _rewrite_workflow_preserve(shape, params):
    if shape is None:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    fmt = _sample_fmt(shape)
    bullets = [b for b in params.get("workflow_bullets", []) if b.strip()]
    foot = params.get("workflow_footnote", "").strip()
    first = _clear_textframe(tf)
    for i, b in enumerate(bullets):
        para = first if i == 0 else tf.add_paragraph()
        _set_number(para)
        _emit(para, b, size_pt=fmt["size_pt"], font=fmt["font"], color=fmt["color"])
    if foot:
        para = tf.add_paragraph(); _no_bullet(para)
        _emit(para, foot, size_pt=11, italic=True, font=fmt["font"], color=fmt["color"])


def _cell_set_preserve(cell, text):
    fmt = _sample_fmt(cell)
    align = cell.text_frame.paragraphs[0].alignment
    para = _clear_textframe(cell.text_frame)
    if align is not None:
        para.alignment = align
    _emit(para, text, size_pt=fmt["size_pt"], bold=fmt["bold"],
          color=fmt["color"], font=fmt["font"])


def _update_table_inplace(table_gf, milestones, total_price):
    table = table_gf.table
    n = len(table.rows)
    for i, m in enumerate(milestones, start=1):
        if i >= n - 1:
            break
        pv = parse_price(m.get("price", ""))
        price_disp = format_money(pv) if pv is not None else str(m.get("price", ""))
        _cell_set_preserve(table.cell(i, 0), m.get("name", ""))
        _cell_set_preserve(table.cell(i, 1), price_disp)
        _cell_set_preserve(table.cell(i, 2), m.get("terms", ""))
    _cell_set_preserve(table.cell(n - 1, 1), total_price)


def _update_boxes_inplace(boxes, milestones):
    for box, m in zip(boxes, milestones):
        tf = box.text_frame
        fmt = _sample_fmt(box)
        label = m.get("short_label", "").strip() or _derive_short_label(m.get("name", ""))
        first = _clear_textframe(tf)
        first.alignment = PP_ALIGN.CENTER
        _emit(first, label, size_pt=fmt["size_pt"], bold=True,
              font=fmt["font"], color=fmt["color"])
        w = parse_weeks(m.get("duration", ""))
        if w:
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            _add_run(p2, fmt_weeks(*w), size_pt=max(7, fmt["size_pt"] - 1),
                     font=fmt["font"], color=fmt["color"])


# ---- read a (possibly hand-edited) deck back into a params dict ----
def _strip_label(text, label):
    t = (text or "").strip()
    low = t.lower()
    if low.startswith(label.lower() + ":"):
        t = t[len(label) + 1:].strip()
    elif low.startswith(label.lower()):
        t = t[len(label):].lstrip(": ").strip()
    return t


def _read_table_milestones(table_gf, boxes):
    table = table_gf.table
    rows = list(table.rows)
    milestones, total = [], ""
    if len(rows) >= 2:
        for i in range(1, len(rows) - 1):
            milestones.append({
                "name": table.cell(i, 0).text_frame.text.strip(),
                "price": table.cell(i, 1).text_frame.text.strip(),
                "terms": table.cell(i, 2).text_frame.text.strip(),
                "short_label": "", "duration": "",
            })
        total = table.cell(len(rows) - 1, 1).text_frame.text.strip()
    # short_label + duration from the banner boxes, matched by position
    for idx, box in enumerate(boxes):
        if idx >= len(milestones):
            break
        paras = [p.text.strip() for p in box.text_frame.paragraphs if p.text.strip()]
        if paras:
            milestones[idx]["short_label"] = paras[0]
        if len(paras) > 1:
            milestones[idx]["duration"] = normalize_weeks_cell(paras[1])
    return milestones, total


def read_deck_into_params(pptx_bytes):
    """Reconstruct a params dict from a generated (and possibly hand-edited) deck.

    Visible text is read from the shapes so manual edits are recovered; non-visible
    fields (transcript, target, ballpark, etc.) come from the embedded state."""
    params = copy.deepcopy(default_params())
    state = _read_state_from_bytes(pptx_bytes)
    if state:
        params.update(state)
    prs = Presentation(io.BytesIO(pptx_bytes))
    slides = list(prs.slides)

    s1 = slides[SLIDE_TITLE]
    sub = _shape_by_id(s1, ID_SUBTITLE)
    if sub is not None:
        nonempty = [p.text.strip() for p in sub.text_frame.paragraphs if p.text.strip()]
        if nonempty:
            params["subtitle"] = nonempty[0]
            params["phase_label"] = nonempty[1] if len(nonempty) > 1 else ""
    date_sh = _shape_by_id(s1, ID_DATE)
    if date_sh is not None and date_sh.text_frame.text.strip():
        params["proposal_date"] = date_sh.text_frame.text.strip()

    s5 = slides[SLIDE_CHALLENGE]
    ch = _shape_by_id(s5, ID_CHALLENGE)
    if ch is not None and ch.text_frame.text.strip():
        params["challenge_text"] = _strip_label(ch.text_frame.text, "Challenge")
    stg = _shape_by_id(s5, ID_STRATEGY)
    if stg is not None and stg.text_frame.text.strip():
        params["strategy_text"] = _strip_label(stg.text_frame.text, "Strategy")

    wf = _shape_by_id(slides[SLIDE_WORKFLOW], ID_WORKFLOW)
    if wf is not None:
        bullets, foot = [], ""
        for p in wf.text_frame.paragraphs:
            t = p.text.strip()
            if not t:
                continue
            if t.startswith("*"):
                foot = t
            else:
                bullets.append(t)
        if bullets:
            params["workflow_bullets"] = bullets
        params["workflow_footnote"] = foot

    s7 = slides[SLIDE_TIMELINE]
    tlt = _shape_by_id(s7, ID_TL_TITLE)
    if tlt is not None and tlt.text_frame.text.strip():
        params["timeline_title"] = tlt.text_frame.text.strip()
    foot = _shape_by_id(s7, ID_TL_FOOT)
    if foot is not None and foot.text_frame.text.strip():
        params["total_note"] = foot.text_frame.text.strip()
    table = _find_table(s7)
    if table is not None:
        ms, total = _read_table_milestones(table, _find_boxes(s7))
        if ms:
            params["milestones"] = ms
        if total:
            params["total_price"] = total
    params["timeline_estimate"] = timeline_line(params.get("milestones", []))
    return params


def revise_proposal_pptx(params, base_pptx_bytes, prev_params=None):
    """Build the next version ONTO the user's edited deck, preserving its formatting.

    Only fields that differ from prev_params are rewritten (so untouched, manually
    formatted shapes are left alone). The milestone table and banner boxes are
    updated in place when the phase count is unchanged, and rebuilt with house
    defaults only when phases were added or removed. Returns new deck bytes with
    refreshed embedded state and an incremented revision counter."""
    prs = Presentation(io.BytesIO(base_pptx_bytes))
    slides = list(prs.slides)
    prev = prev_params or {}

    def changed(key):
        return prev.get(key) != params.get(key)

    s1 = slides[SLIDE_TITLE]
    if changed("subtitle") or changed("phase_label"):
        _rewrite_subtitle_preserve(_shape_by_id(s1, ID_SUBTITLE), params)
    if changed("proposal_date"):
        _set_text_preserve(_shape_by_id(s1, ID_DATE), params.get("proposal_date", ""))

    s5 = slides[SLIDE_CHALLENGE]
    if changed("challenge_text"):
        _set_label_body_preserve(_shape_by_id(s5, ID_CHALLENGE), "Challenge:",
                                 params.get("challenge_text", ""))
    if changed("strategy_text"):
        _set_label_body_preserve(_shape_by_id(s5, ID_STRATEGY), "Strategy:",
                                 params.get("strategy_text", ""))

    if changed("workflow_bullets") or changed("workflow_footnote"):
        _rewrite_workflow_preserve(_shape_by_id(slides[SLIDE_WORKFLOW], ID_WORKFLOW), params)

    s7 = slides[SLIDE_TIMELINE]
    if changed("timeline_title"):
        _set_text_preserve(_shape_by_id(s7, ID_TL_TITLE), params.get("timeline_title", ""))
    # timeline line is derived from durations -> always refresh it
    _set_text_preserve(_shape_by_id(s7, ID_TL_LINE), params.get("timeline_estimate", ""))
    if changed("total_note"):
        _set_text_preserve(_shape_by_id(s7, ID_TL_FOOT), params.get("total_note", ""))

    new_ms = params.get("milestones", [])
    table = _find_table(s7)
    boxes = _find_boxes(s7)
    cur_rows = (len(table.table.rows) - 2) if table is not None else -1
    count_same = (table is not None and len(new_ms) == cur_rows and len(boxes) == len(new_ms))
    if new_ms:
        if count_same:
            if changed("milestones") or changed("total_price"):
                _update_table_inplace(table, new_ms, params.get("total_price", ""))
            if changed("milestones"):
                _update_boxes_inplace(boxes, new_ms)
        else:
            tb = _rebuild_milestones_table(s7, params)   # rebuild with house defaults
            _rebuild_timeline_boxes(s7, new_ms)
            ft = _shape_by_id(s7, ID_TL_FOOT)
            if ft is not None and tb is not None:
                ft.top = Emu(int(tb) + TABLE_FOOT_GAP_EMU)

    state = dict(params)
    state["_revision"] = int(prev.get("_revision", params.get("_revision", 0)) or 0) + 1
    bio = io.BytesIO()
    prs.save(bio)
    return _embed_state_in_bytes(bio.getvalue(), state)


# --------------------------------------------------------------------------- #
# Defaults (the Macoska / UMass Boston example, used as a worked demo)
# --------------------------------------------------------------------------- #
def default_params():
    return {
        "customer_short": "UMass Boston",
        "customer_full": "University of Massachusetts Boston",
        "pi_contact": "Dr. Jill Macoska",
        "institution_type": "Academic",
        "proposal_date": datetime.now().strftime("%B %d, %Y"),
        "subtitle": ("Confidential Proposal to UMass Boston for a Urinary "
                     "Collagen Aptamer Diagnostic"),
        "phase_label": "(Phase I \u2013 understanding existing aptamer suitability in urine)",
        # Target & strategy parameters
        "target": "fibrillar collagen (Gly-X-Y tripeptide repeat)",
        "target_type": "Protein",
        "existing_aptamer": True,
        "existing_aptamer_desc": ("an aptamer raised against the Gly-X-Y tripeptide "
                                  "repeat common to the fibrillar collagens"),
        "biological_matrix": "urine (collagen-free / synthetic urine for spike-in)",
        "off_targets": "urine background constituents",
        "assay_format_goal": "point-of-care (electrochemical or lateral-flow) assay",
        "kd_method": KD_METHOD_BY_TYPE["Protein"],
        "phases_included": "Phase 1 + Phase 2 (if needed)",
        "background_problem": ("Dr. Macoska's lab has shown that fibrillar collagen "
                               "accumulates in the urine of men with lower urinary "
                               "tract dysfunction while controls show almost none; no "
                               "non-invasive test for urinary tract fibrosis exists, "
                               "and the dye-based Sircol/Sirius Red assay cannot be "
                               "licensed for clinical use."),
        # Drafted slide prose
        "challenge_text": (
            "Dr. Macoska\u2019s lab has shown that fibrillar collagen \u2013 the hallmark "
            "of tissue fibrosis \u2013 accumulates in the urine of men with lower urinary "
            "tract dysfunction, while controls show almost none. No non-invasive test for "
            "urinary tract fibrosis exists today, and the dye-based assay used to date "
            "(Sircol / Sirius Red) cannot be licensed for clinical use. UMass Boston seeks "
            "a point-of-care reagent to quantify fibrillar collagen in urine."),
        "strategy_text": (
            "Base Pair already holds an aptamer raised against the Gly-X-Y tripeptide "
            "repeat common to the fibrillar collagens. Because aptamer binding can be "
            "matrix-dependent, the first objective is to confirm that this existing "
            "aptamer recognizes fibrillar collagen in urine and to determine its binding "
            "affinity (KD) in that matrix, using collagen spiked into collagen-free / "
            "synthetic urine to establish a quantitative dose-response. If the existing "
            "aptamer binds with suitable affinity, Base Pair will advance it toward a "
            "quantitative point-of-care assay. If binding in urine is inadequate, Base "
            "Pair will perform a new SELEX campaign directly in a urine matrix, with "
            "counter-selection against off-targets defined with UMass Boston."),
        "workflow_bullets": [
            "Obtain the existing anti-collagen aptamer, a purified fibrillar-collagen reagent, and collagen-free / synthetic urine",
            "Phase 1 \u2013 Test the existing aptamer for binding to fibrillar collagen in urine and determine its binding affinity (KD) using collagen spiked into urine",
            "Establish a quantitative dose-response across spiked collagen concentrations",
            "Confirm specificity against urine background and assess assay window vs. clinical collagen levels",
            "Phase 2 (if needed) \u2013 Perform a new SELEX campaign directly in a urine matrix to generate aptamers optimized for urine, with counter-selection against off-targets",
            "Screen and rank candidates by NGS; complete KD determination of best candidates via MicroScale Thermophoresis (MST) or another appropriate method",
            "Synthesize sample aptamer materials for point-of-care assay development by UMass Boston (electrochemical or lateral-flow format)*",
        ],
        "workflow_footnote": "* Phase 2 SELEX and any large-scale synthesis are scoped and quoted separately.",
        # Timeline & milestones
        "timeline_title": "Timeline and Milestones \u2013 Phase 1 Feasibility",
        "timeline_estimate": "Timeline:       ~5\u20138 weeks from receipt of materials",
        "milestones": [
            {"name": "Project Initiation & Materials", "price": "$2,500",
             "terms": "Due upon execution of the Aptamer Development Agreement.",
             "short_label": "Obtain materials", "duration": "1"},
            {"name": "Aptamer Binding Test in Urine (yes/no ranging expt)", "price": "$1,500",
             "terms": "Due within 30 days of Base Pair reporting binding results.",
             "short_label": "Binding test in urine", "duration": "1-2"},
            {"name": "KD Determination & Dose-Response in Urine", "price": "$4,500",
             "terms": "Due within 30 days of delivery of the Binding / KD report.",
             "short_label": "KD & dose-response", "duration": "2-3"},
            {"name": "Sample Aptamer Synthesis for Assay Development", "price": "$1,400",
             "terms": "Due within 30 days of shipment of sample materials.",
             "short_label": "Sample synthesis", "duration": "1-2"},
        ],
        "total_price": "$9,900*",
        "total_note": ("* Final pricing provided in a formal quote. Phase 2 SELEX in urine "
                       "is scoped separately if needed."),
    }


# --------------------------------------------------------------------------- #
# LLM layer (swap provider here only)
# --------------------------------------------------------------------------- #
def call_llm(client, model, system, user, max_tokens=1600):
    """Single choke-point for the LLM (Anthropic / Claude). Returns assistant text."""
    resp = client.messages.create(
        model=model,
        max_tokens=max_tokens,
        temperature=0.2,
        system=system,
        messages=[{"role": "user", "content": user}],
    )
    return "".join(b.text for b in resp.content if getattr(b, "type", "") == "text").strip()


def _strip_fences(text):
    t = text.strip()
    if t.startswith("```"):
        t = t.split("```", 2)[1] if t.count("```") >= 2 else t.strip("`")
        if t.lstrip().lower().startswith("json"):
            t = t.lstrip()[4:]
    return t.strip()


EXTRACT_SYSTEM = """You are an analyst at Base Pair Biotechnologies, an aptamer
discovery company. You read a raw, messy transcript of a discovery call with a
prospective customer and extract the facts needed to draft a commercial proposal.

Return ONLY a JSON object (no prose, no markdown fences) with these keys:
  customer_short        short customer/institution name (e.g. "UMass Boston")
  customer_full         full institution name if stated, else same as short
  pi_contact            the customer's lead scientist / PI, with title if given
  institution_type      one of: Academic, Biotech, Pharma, Startup, Government, Other
  target                the molecule/entity the aptamer must bind (be specific;
                        include any epitope/motif mentioned)
  target_type           one of: Protein, Peptide epitope, Small molecule,
                        Whole cells, Viral particle, Other
  existing_aptamer      true if Base Pair already holds a relevant aptamer, else false
  existing_aptamer_desc short description of that existing aptamer, else ""
  biological_matrix     the sample type the assay must work in (e.g. "urine")
  off_targets           anything that must be avoided / counter-selected, else ""
  assay_format_goal     the end assay the customer wants (e.g. "point-of-care
                        lateral-flow"), else ""
  background_problem    2-4 sentences on the customer's scientific/clinical problem,
                        the gap in current methods, and what they want, in neutral
                        third-person voice
  subtitle              a one-line title-slide subtitle in the exact form
                        "Confidential Proposal to <customer_short> for a <descriptor>
                        Aptamer Diagnostic" (or "...Aptamer Development Program" if it
                        is not a diagnostic). <descriptor> names the matrix/target,
                        e.g. "Urinary Collagen", "Serum Nattokinase".
  phase_label           a short parenthetical scope note for the title slide, e.g.
                        "(Phase I - feasibility of an existing aptamer in serum)" or
                        "(Phase I - de novo aptamer discovery)". Reflect whether an
                        existing aptamer is being tested vs. a new SELEX campaign.
                        Leave "" if unclear.
Leave any unknown string blank ("") rather than guessing. Do not invent pricing,
dates, or numbers that were not stated."""

DRAFT_SYSTEM = """You write the project-specific slides of a Base Pair Biotechnologies
commercial proposal. Base Pair discovers custom DNA/RNA aptamers. The voice is
professional, confident, client-facing, and concrete -- no marketing fluff, no
hyperbole, present/future tense.

You are given a JSON object of project parameters. Return ONLY a JSON object (no
markdown fences) with these keys:
  challenge_text   ONE paragraph (no leading "Challenge:" label) describing the
                   customer's problem and what they need -- 3-5 sentences.
  strategy_text    ONE paragraph (no leading "Strategy:" label) describing Base
                   Pair's approach: if an existing aptamer is available, the first
                   objective is to confirm it recognizes the target in the stated
                   matrix and measure KD there (using a spiked dose-response in a
                   blank/synthetic matrix); state the success branch (advance toward
                   the desired assay format) and the fallback branch (a new SELEX
                   campaign performed directly in the matrix, with counter-selection
                   against the stated off-targets). If no existing aptamer, lead with
                   a new SELEX campaign in the matrix. 4-6 sentences.
  workflow_bullets a JSON array of 6-8 short imperative bullet strings describing the
                   stepwise plan (obtain materials; Phase 1 binding test in matrix;
                   KD/dose-response; specificity vs background; Phase 2 new SELEX if
                   needed with counter-selection; NGS screen + KD by the appropriate
                   method; synthesize sample materials for the customer's assay).
                   Use the customer's actual target, matrix, off-targets, KD method,
                   and assay format. End the synthesis bullet with a "*".
  phases           a JSON array of 4-6 objects, each {"short_label": "...",
                   "duration": "...", "name": "..."} representing the project's
                   payment/timeline phases. short_label is 2-4 words for a timeline
                   banner box; duration is the phase length in WEEKS as a bare number
                   or range only (e.g. "1" or "2-3") -- no words, no "~", no "weeks";
                   name is the fuller milestone description for a table.
                   These should summarize the workflow into billing phases (typically:
                   initiation/materials; the key binding test; KD/dose-response;
                   sample synthesis; optionally a Phase 2 SELEX phase). DO NOT include
                   any price or payment terms -- the user sets those.
Keep numbers/affinity methods consistent with the parameters. Do not include pricing
or payment terms anywhere."""


def llm_extract(client, model, transcript):
    user = "TRANSCRIPT:\n\n" + transcript[:24000]
    raw = call_llm(client, model, EXTRACT_SYSTEM, user, max_tokens=1200)
    return json.loads(_strip_fences(raw))


def merge_phases_keep_pricing(milestones, phases):
    """Merge LLM-suggested phases into milestones by row position, preserving any
    price/terms already entered. Pricing is never LLM-generated."""
    old = milestones or []
    merged = []
    for i, item in enumerate(list(phases)[:6]):
        prev = old[i] if i < len(old) else {}
        merged.append({
            "short_label": str(item.get("short_label", "") or ""),
            "duration": normalize_weeks_cell(item.get("duration", "")),
            "name": str(item.get("name", "") or ""),
            "price": prev.get("price", ""),
            "terms": prev.get("terms", ""),
        })
    return merged


def llm_draft(client, model, params):
    keys = ["customer_short", "pi_contact", "target", "target_type",
            "existing_aptamer", "existing_aptamer_desc", "biological_matrix",
            "off_targets", "assay_format_goal", "kd_method", "phases_included",
            "background_problem"]
    payload = {k: params.get(k) for k in keys}
    user = "PARAMETERS:\n\n" + json.dumps(payload, indent=2)
    raw = call_llm(client, model, DRAFT_SYSTEM, user, max_tokens=1600)
    return json.loads(_strip_fences(raw))


REVISE_SYSTEM = """You revise an existing Base Pair Biotechnologies commercial
proposal in light of customer feedback. Base Pair discovers custom DNA/RNA aptamers.
Voice: professional, confident, client-facing, concrete. Present/future tense.

You are given THREE things: (1) the original discovery-call transcript, (2) the
CURRENT proposal (its parameters and drafted slide text), and (3) new FEEDBACK
(an email thread, notes, or a follow-up call). Produce a REVISED proposal that
honors the feedback, preserves anything the feedback does not change, and reflects
any agreed change in strategy, sequencing, or scope.

Return ONLY a JSON object (no markdown fences) with these keys (include a key only
if it should change; omitted keys keep their current value):
  subtitle           revised title-slide subtitle in the form "Confidential Proposal
                     to <customer> for a <descriptor> Aptamer <Diagnostic|Development
                     Program>" if the project framing changed; else omit.
  phase_label        revised parenthetical scope note for the title slide if scope/
                     phasing changed (e.g. "(Phase I - denatured-form assay)"); else omit.
  challenge_text     ONE paragraph (no "Challenge:" label), 3-5 sentences, updated for
                     the feedback.
  strategy_text      ONE paragraph (no "Strategy:" label), 4-6 sentences, updated for
                     the feedback (e.g. revised order of work, new sub-targets,
                     deferred phases, changed matrix/denaturation).
  workflow_bullets   a JSON array of 6-8 short imperative bullets reflecting the
                     revised plan. End the synthesis/deliverable bullet with "*".
  phases             a JSON array of 4-6 objects {"short_label","duration","name"}
                     for the revised payment/timeline phases. You MAY restructure
                     these (reorder, split, add, or remove phases) to match the
                     feedback -- e.g. sequencing a denatured-form phase before a
                     harder neat-serum phase. duration is WEEKS as a bare number or
                     range only (e.g. "1" or "2-3"). DO NOT include price or terms.
ABSOLUTE RULE: never output prices, payment terms, or dollar amounts -- pricing is
set by the user only. Keep affinity methods/numbers consistent unless the feedback
changes them."""


def llm_revise(client, model, params, transcript, feedback):
    keys = ["customer_short", "pi_contact", "target", "target_type",
            "existing_aptamer", "existing_aptamer_desc", "biological_matrix",
            "off_targets", "assay_format_goal", "kd_method", "phases_included",
            "background_problem", "subtitle", "phase_label", "challenge_text",
            "strategy_text", "workflow_bullets"]
    current = {k: params.get(k) for k in keys}
    current["phases"] = [{"short_label": m.get("short_label", ""),
                          "duration": m.get("duration", ""),
                          "name": m.get("name", "")}
                         for m in params.get("milestones", [])]
    user = ("ORIGINAL TRANSCRIPT:\n\n" + (transcript or "")[:16000]
            + "\n\n----\nCURRENT PROPOSAL:\n\n" + json.dumps(current, indent=2)
            + "\n\n----\nFEEDBACK:\n\n" + (feedback or "")[:8000])
    raw = call_llm(client, model, REVISE_SYSTEM, user, max_tokens=1800)
    return json.loads(_strip_fences(raw))


# --------------------------------------------------------------------------- #
# Streamlit UI
# --------------------------------------------------------------------------- #
def run_app():
    import streamlit as st

    st.set_page_config(page_title="Base Pair Proposal Generator", layout="wide")

    # ----- BEGIN BasePair Secure Password Gate -----
    # Password is read from Streamlit secrets / env var / local file -- never
    # hardcoded -- so this repository can safely be public.
    _app_pw = load_app_password()
    if not _app_pw:
        st.error("No app password is configured, so access is locked. Set "
                 "**APP_PASSWORD** in Streamlit \u2192 Settings \u2192 Secrets (or the "
                 "BPB_APP_PASSWORD environment variable when running locally) to "
                 "enable the app.")
        st.stop()
    password = st.text_input("Enter password (same as BasePair wifi password):",
                             type="password")
    if password != _app_pw:
        st.warning("Incorrect password. Please try again.")
        st.stop()
    # ----- END Password Gate -----

    # logo is optional; the deck carries its own
    if os.path.exists(LOGO_PATH):
        st.image(LOGO_PATH, width=150)
    st.title("Base Pair \u2014 Commercial Proposal Generator")
    st.caption("Otter transcript \u2192 parameters \u2192 drafted slides \u2192 proposal deck (.pptx). "
               "Boilerplate tech slides are preserved; only project-specific slides are written.")

    with st.expander("About This App"):
        st.write(
            "Turns an Otter.ai discovery-call transcript into a populated, non-binding "
            "commercial proposal deck (.pptx) in Base Pair's standard format. The company "
            "technology slides are preserved unchanged; only the project-specific slides "
            "(title, Challenge & Strategy, Workflow, Timeline & Milestones) are written.")
        st.markdown(
            "**How to use:**\n"
            "1. Paste or upload the Otter transcript in section 1, then click "
            "**Draft Challenge & Strategy** \u2014 Claude reads the call and drafts the "
            "title, Challenge, Strategy, Workflow steps, and project phases.\n"
            "2. Review and edit everything in the sidebar and section 2 (all text is editable).\n"
            "3. In section 3, set phase durations (in weeks) and enter pricing \u2014 either "
            "type each milestone fee or enter a **Ballpark project total** and click "
            "**Approximate prices** (Project Initiation is frontloaded at 31%; the rest is "
            "distributed by phase duration). The timeline line is computed from the durations.\n"
            "4. Click **Build .pptx** in section 4 to download the deck. You can refine the "
            "deck further in PowerPoint afterward.\n\n"
            "Pricing is never set by the LLM \u2014 you always enter or approve it.")
        st.markdown(
            "**A note on \u201cRecent proposals\u201d (sidebar):** *Save current to history* keeps "
            "a snapshot of the form so you can reload or branch from it \u2014 but only for "
            "your current session. It is **not** a saved library: it clears when you close "
            "the tab, when the session times out, or whenever the app is updated/redeployed, "
            "and it is not shared between users. To keep a proposal, download the **.pptx** "
            "(via Build) \u2014 that file is the permanent record.")

    # Pulse/glow animation to draw the eye to the two key action buttons until
    # each is used. We give each button a stable key; Streamlit renders that as
    # a container class `st-key-<key>`, which is a reliable styling hook.
    _draft_done = st.session_state.get("draft_done", False)
    _price_done = st.session_state.get("price_done", False)
    _targets = []
    if not _draft_done:
        _targets.append(".st-key-btn_draft button")
    if not _price_done:
        _targets.append(".st-key-btn_price button")
    if _targets:
        sel = ", ".join(_targets)
        st.markdown(
            "<style>\n"
            "@keyframes bpbPulse {\n"
            "  0%   { box-shadow: 0 0 0 0 rgba(84,130,53,0.6); }\n"
            "  70%  { box-shadow: 0 0 0 14px rgba(84,130,53,0); }\n"
            "  100% { box-shadow: 0 0 0 0 rgba(84,130,53,0); }\n"
            "}\n"
            f"{sel} {{\n"
            "  animation: bpbPulse 1.5s infinite !important;\n"
            "  border: 2px solid #548235 !important;\n"
            "}\n"
            "</style>", unsafe_allow_html=True)

    # ---- session init ----
    if "params" not in st.session_state:
        st.session_state["params"] = default_params()
    if "history" not in st.session_state:
        st.session_state["history"] = []
    st.session_state.setdefault("ms_editor_nonce", 0)

    def _bump_ms_editor():
        # Re-mount the milestone editor so it reloads from params (used whenever
        # we set prices/milestones programmatically, e.g. Approximate prices,
        # deck load, draft/revise). Without this the editor keeps its old cells
        # and clobbers the new values on the next rerun.
        st.session_state["ms_editor_nonce"] = st.session_state.get("ms_editor_nonce", 0) + 1

    P = st.session_state["params"]

    # ---- Anthropic client (key persisted on this machine) ----
    if "anthropic_api_key" not in st.session_state:
        st.session_state["anthropic_api_key"] = load_saved_key()
    with st.sidebar:
        st.header("LLM")
        if st.session_state["anthropic_api_key"]:
            st.success("API key loaded \u2713")
            if st.button("Forget saved key"):
                forget_key()
                st.session_state["anthropic_api_key"] = ""
                st.rerun()
        else:
            entered = st.text_input(
                "Anthropic API Key", type="password",
                help="Paste once \u2014 it's saved on this computer so you won't be asked again.")
            remember = st.checkbox("Remember on this computer", value=True)
            if entered:
                st.session_state["anthropic_api_key"] = entered.strip()
                if remember and save_key(entered):
                    st.success("Saved \u2014 you won't need to paste it again.")
                elif remember:
                    st.warning("Couldn't write the key file; it will persist for this session only.")
                st.rerun()
        model = st.text_input("Model", value=st.session_state.get("model", "claude-sonnet-4-6"))
        st.session_state["model"] = model

    client = None
    if st.session_state["anthropic_api_key"]:
        try:
            import anthropic
            client = anthropic.Anthropic(api_key=st.session_state["anthropic_api_key"])
        except Exception as e:
            st.sidebar.error(f"Anthropic client error: {e}")

    # ---- recent proposals ----
    with st.sidebar:
        st.header("Recent proposals")
        if st.button("Save current to history"):
            st.session_state["history"].insert(0, {
                "label": f"{P['customer_short']} \u2013 {P['target'][:24]} "
                         f"\u2013 {datetime.now().strftime('%Y.%m.%d %H:%M')}",
                "params": copy.deepcopy(P)})
            st.session_state["history"] = st.session_state["history"][:6]
            st.success("Saved.")
        for i, h in enumerate(st.session_state["history"]):
            if st.button(h["label"], key=f"hist_{i}"):
                st.session_state["params"] = copy.deepcopy(h["params"])
                _bump_ms_editor()
                st.rerun()

    # ====================================================================== #
    # STEP 1 — Transcript / revision input
    # ====================================================================== #
    st.subheader("1) Discovery-call transcript")
    mode = st.radio(
        "Mode", ["New proposal", "Revise existing"], horizontal=True,
        help="New: draft a fresh proposal from a transcript. Revise: apply "
             "customer feedback (e.g. an email thread) to the proposal currently "
             "loaded below, keeping your prices.")
    up = st.file_uploader("Upload Otter transcript (.txt)", type=["txt"])
    transcript = ""
    if up is not None:
        transcript = up.read().decode("utf-8", errors="ignore")
    transcript = st.text_area("…or paste transcript here", value=transcript, height=160)

    feedback = ""
    if mode == "Revise existing":
        st.markdown(
            "**Revise from your edited deck.** Upload the **.pptx you downloaded and "
            "tweaked in PowerPoint** \u2014 your manual formatting and text edits are kept, "
            "and only what the feedback changes is rewritten. You can repeat this for "
            "multiple rounds (each downloaded revision can be re-uploaded here).")
        deck_up = st.file_uploader(
            "Upload the edited proposal deck (.pptx) to revise", type=["pptx"],
            key="deck_up")
        if deck_up is not None:
            sig = (deck_up.name, deck_up.size)
            if st.session_state.get("loaded_deck_sig") != sig:
                try:
                    raw = deck_up.getvalue()
                    base = read_deck_into_params(raw)
                    st.session_state["base_deck_bytes"] = raw
                    st.session_state["base_deck_params"] = copy.deepcopy(base)
                    st.session_state["params"] = base
                    st.session_state["loaded_deck_sig"] = sig
                    _bump_ms_editor()
                    P = st.session_state["params"]
                    rv = (_read_state_from_bytes(raw) or {}).get("_revision", 0)
                    st.success(f"Loaded the deck (revision {rv}). Its current text is now in "
                               f"the fields below. Paste the feedback, then "
                               f"**Apply feedback & re-draft** \u2014 or just tweak pricing and "
                               f"Build.")
                except Exception as e:
                    st.error(f"Couldn't read that deck: {e}")
            else:
                st.caption("Deck loaded. Edits you make below are kept; re-upload a "
                           "different file to start from another deck.")
        else:
            st.session_state["loaded_deck_sig"] = None
            if st.session_state.get("base_deck_bytes"):
                st.caption("A deck is loaded for revision. Upload a different one to replace it.")
        feedback = st.text_area(
            "Follow-up / feedback (paste the email thread, notes, or a follow-up call)",
            height=160)

    c1, c2 = st.columns([1, 3])
    with c1:
        if mode == "New proposal":
            do_action = st.button("Draft Challenge & Strategy", type="primary",
                                  key="btn_draft")
        else:
            do_action = st.button("Apply feedback & re-draft", type="primary",
                                  key="btn_draft")
    with c2:
        if mode == "New proposal":
            st.caption("Reads the transcript, then drafts the title, Challenge, "
                       "Strategy, Workflow & phases. Pricing is never set "
                       "automatically \u2014 you enter that in section 3.")
        else:
            st.caption("Applies the feedback to the proposal currently loaded below "
                       "(it can reorder/split phases). Your prices are preserved. "
                       "Review everything before Build.")

    if do_action:
        if not client:
            st.error("Provide an Anthropic API key first.")
        elif mode == "New proposal" and not transcript.strip():
            st.error("Add a transcript first.")
        elif mode == "Revise existing" and not feedback.strip():
            st.error("Paste the feedback to apply first.")
        else:
            try:
                if mode == "New proposal":
                    st.session_state["base_deck_bytes"] = None
                    st.session_state["base_deck_params"] = None
                    st.session_state["loaded_deck_sig"] = None
                    with st.spinner("Reading the call…"):
                        extracted = llm_extract(client, model, transcript)
                        for k, v in extracted.items():
                            if v not in (None, ""):
                                P[k] = v
                        tt = P.get("target_type", "Protein")
                        P["kd_method"] = KD_METHOD_BY_TYPE.get(tt, KD_METHOD_BY_TYPE["Other"])
                        P["proposal_date"] = datetime.now().strftime("%B %d, %Y")
                        P["_transcript"] = transcript
                    with st.spinner("Drafting in Base Pair voice…"):
                        drafted = llm_draft(client, model, P)
                        P["challenge_text"] = drafted.get("challenge_text", P["challenge_text"])
                        P["strategy_text"] = drafted.get("strategy_text", P["strategy_text"])
                        wb = drafted.get("workflow_bullets")
                        if isinstance(wb, list) and wb:
                            P["workflow_bullets"] = wb
                        ph = drafted.get("phases")
                        if isinstance(ph, list) and ph:
                            P["milestones"] = merge_phases_keep_pricing(P.get("milestones", []), ph)
                    st.session_state["params"] = P
                    st.session_state["draft_done"] = True
                    _bump_ms_editor()
                    st.success("Drafted from the transcript. Review and edit every "
                               "section below \u2014 then enter pricing in section 3 and Build.")
                else:  # Revise existing
                    on_deck = bool(st.session_state.get("base_deck_bytes"))
                    src_transcript = (P.get("_transcript") or transcript) if on_deck else transcript
                    with st.spinner("Applying feedback…"):
                        rev = llm_revise(client, model, P, src_transcript, feedback)
                        for k in ("subtitle", "phase_label", "challenge_text",
                                  "strategy_text"):
                            if rev.get(k):
                                P[k] = rev[k]
                        wb = rev.get("workflow_bullets")
                        if isinstance(wb, list) and wb:
                            P["workflow_bullets"] = wb
                        ph = rev.get("phases")
                        if isinstance(ph, list) and ph:
                            P["milestones"] = merge_phases_keep_pricing(P.get("milestones", []), ph)
                        if not on_deck:
                            P["proposal_date"] = datetime.now().strftime("%B %d, %Y")
                    st.session_state["params"] = P
                    _bump_ms_editor()
                    msg = ("Revised per the feedback. Review the changes below "
                           "(phases may have been restructured) \u2014 your prices were kept.")
                    if on_deck:
                        msg += (" When you Build, the changes are written **onto your "
                                "uploaded deck**, so your manual formatting is preserved.")
                    st.success(msg)
            except Exception as e:
                st.error(f"Generation failed: {e}")

    # ====================================================================== #
    # Sidebar — parameter review/edit
    # ====================================================================== #
    with st.sidebar:
        st.header("Parameters")
        with st.expander("Customer & contacts", expanded=True):
            P["customer_short"] = st.text_input("Customer (short)", P["customer_short"])
            P["customer_full"] = st.text_input("Customer (full)", P.get("customer_full", ""))
            P["pi_contact"] = st.text_input("Lead contact / PI", P.get("pi_contact", ""))
            P["institution_type"] = st.selectbox(
                "Institution type",
                ["Academic", "Biotech", "Pharma", "Startup", "Government", "Other"],
                index=max(0, ["Academic", "Biotech", "Pharma", "Startup",
                              "Government", "Other"].index(P.get("institution_type", "Academic"))
                          if P.get("institution_type") in
                          ["Academic", "Biotech", "Pharma", "Startup", "Government", "Other"] else 0))
            P["proposal_date"] = st.text_input("Proposal date", P["proposal_date"])

        with st.expander("Target & strategy", expanded=True):
            P["target"] = st.text_input("Target", P.get("target", ""))
            prev_tt = P.get("target_type", "Protein")
            tt = st.selectbox("Target type", TARGET_TYPES,
                              index=TARGET_TYPES.index(prev_tt)
                              if prev_tt in TARGET_TYPES else 0)
            # changing the target type re-suggests the affinity method
            if tt != prev_tt:
                P["kd_method"] = KD_METHOD_BY_TYPE.get(tt, KD_METHOD_BLI)
            P["target_type"] = tt
            P["existing_aptamer"] = st.checkbox(
                "Base Pair already holds a relevant aptamer",
                value=bool(P.get("existing_aptamer", False)))
            P["existing_aptamer_desc"] = st.text_area(
                "Existing aptamer (description)", P.get("existing_aptamer_desc", ""))
            P["biological_matrix"] = st.text_input("Biological matrix", P.get("biological_matrix", ""))
            P["off_targets"] = st.text_input("Off-targets to avoid (counter-selection)",
                                             P.get("off_targets", ""))
            P["assay_format_goal"] = st.text_input("Assay format goal", P.get("assay_format_goal", ""))
            cur_kd = P.get("kd_method") or KD_METHOD_BY_TYPE.get(tt, KD_METHOD_BLI)
            kd_opts = KD_METHOD_OPTIONS if cur_kd in KD_METHOD_OPTIONS \
                else KD_METHOD_OPTIONS + [cur_kd]
            P["kd_method"] = st.selectbox(
                "Affinity ($K_D$) determination method", kd_opts,
                index=kd_opts.index(cur_kd),
                help="Auto-suggested from Target type; override here if needed.")
            P["phases_included"] = st.selectbox(
                "Phases", ["Phase 1 only", "Phase 1 + Phase 2 (if needed)"],
                index=1 if "Phase 2" in P.get("phases_included", "") else 0)

    # ====================================================================== #
    # STEP 2 — Review & edit the drafted slides
    # ====================================================================== #
    st.subheader("2) Review & edit the drafted slides")
    P["background_problem"] = st.text_area(
        "Background / problem \u2014 context the Challenge was drafted from (not placed on a slide verbatim)",
        P.get("background_problem", ""), height=120)

    st.markdown("**Slide 1 \u2014 title page** (text placed verbatim; not LLM-drafted)")
    P["subtitle"] = st.text_area("Subtitle line", P.get("subtitle", ""), height=70)
    P["phase_label"] = st.text_input("Phase label (italic)", P.get("phase_label", ""))

    st.markdown("**Slide 5 \u2014 Challenge & Strategy**")
    P["challenge_text"] = st.text_area("Challenge", P["challenge_text"], height=120)
    P["strategy_text"] = st.text_area("Strategy", P["strategy_text"], height=160)

    st.markdown("**Slide 6 \u2014 Workflow steps** — one field per bullet, in order")
    try:
        import pandas as pd
        wdf = pd.DataFrame({"Workflow step": P["workflow_bullets"] or [""]})
        wedit = st.data_editor(wdf, num_rows="dynamic", use_container_width=True,
                               key="workflow_editor",
                               column_config={"Workflow step": st.column_config.TextColumn(
                                   "Workflow step", width="large")})
        P["workflow_bullets"] = [str(x) for x in wedit["Workflow step"].tolist()
                                 if str(x).strip() and str(x).lower() != "nan"]
    except Exception:
        # fallback if pandas/data_editor unavailable
        n = st.number_input("Number of workflow steps", min_value=1, max_value=20,
                            value=max(1, len(P["workflow_bullets"])), key="wf_n")
        new = []
        for i in range(int(n)):
            cur = P["workflow_bullets"][i] if i < len(P["workflow_bullets"]) else ""
            new.append(st.text_input(f"Step {i+1}", value=cur, key=f"wfb_{i}"))
        P["workflow_bullets"] = [b for b in new if b.strip()]
    P["workflow_footnote"] = st.text_input("Workflow footnote", P.get("workflow_footnote", ""))

    # ====================================================================== #
    # STEP 3 — Timeline & milestones (manual)
    # ====================================================================== #
    st.subheader("3) Project phases, timeline & milestones (you enter pricing)")
    P["timeline_title"] = st.text_input("Slide 7 title", P["timeline_title"])

    st.markdown(
        "**Phases / payment milestones** — each row is one phase: it becomes both a "
        "**timeline banner box** (Phase label + Duration) and a **milestone table row** "
        "(Milestone + Price + Terms). Durations are in **weeks** (number or range, e.g. "
        "`2` or `2-3`). Add or remove rows to fit the project (**4\u20136 recommended**). "
        "Price auto-formats to **$#,###** (no cents).")
    try:
        import pandas as pd
        cols = ["short_label", "duration", "name", "price", "terms"]
        rows = [{c: m.get(c, "") for c in cols} for m in P["milestones"]]
        df = pd.DataFrame(rows, columns=cols)
        if df.empty:
            df = pd.DataFrame([{c: "" for c in cols}])
        edited = st.data_editor(
            df, num_rows="dynamic", use_container_width=True,
            key=f"ms_editor_{st.session_state['ms_editor_nonce']}",
            column_config={
                "short_label": st.column_config.TextColumn(
                    "Phase (banner)", help="Short label for the timeline box; "
                    "blank = auto from Milestone name.", width="small"),
                "duration": st.column_config.TextColumn(
                    "Duration (weeks)", help="Weeks only \u2014 a number or range, e.g. 2 or 2-3. "
                    "Words are stripped automatically.", width="small"),
                "name": st.column_config.TextColumn("Milestone", width="medium"),
                "price": st.column_config.TextColumn(
                    "Price (USD)", help="Type a number; shows as $#,### with no cents."),
                "terms": st.column_config.TextColumn("Payment Terms", width="large"),
            })

        def _fmt_price_cell(v):
            pv = parse_price(v)
            if pv is not None:
                return format_money(pv)
            return "" if v is None else str(v)

        P["milestones"] = [
            {"short_label": str(r.get("short_label", "") or ""),
             "duration": normalize_weeks_cell(r.get("duration", "")),
             "name": str(r.get("name", "") or ""),
             "price": _fmt_price_cell(r.get("price", "")),
             "terms": str(r.get("terms", "") or "")}
            for _, r in edited.iterrows()
            if str(r.get("name", "") or "").strip()
            or str(r.get("short_label", "") or "").strip()
            or str(r.get("price", "") or "").strip()]
    except Exception:
        st.info("Install pandas for the milestone editor; using raw fields.")

    n_ms = len(P["milestones"])
    if n_ms and not (4 <= n_ms <= 6):
        st.warning(f"You have {n_ms} phase(s). 4\u20136 is recommended for a balanced "
                   f"slide 7 banner and milestone table (it will still generate).")

    # ---- computed timeline line (from the Duration column) ----
    P["timeline_estimate"] = timeline_line(P["milestones"])
    tl = compute_timeline(P["milestones"])
    st.markdown(f"**Timeline (auto from durations):** {fmt_weeks(*tl)} from receipt of "
                f"materials" if tl else "**Timeline:** add phase durations (weeks) above.")

    # ---- ballpark pricing ----
    st.markdown("**Ballpark pricing** — enter a target project total and approximate the "
                "per-milestone fees. **Project Initiation (first row) is frontloaded at "
                "31%**; the rest is split across the other milestones by the chosen "
                "method. Then fine-tune any row above.")
    bc1, bc2, bc3 = st.columns([1.2, 1.2, 1])
    with bc1:
        ballpark = st.number_input("Ballpark project total (USD)", min_value=0,
                                   step=500, value=int(P.get("ballpark", 0) or 0))
        P["ballpark"] = ballpark
    with bc2:
        method_label = st.selectbox("Distribute across milestones by",
                                    ["Duration (weeks)", "Even split"])
    with bc3:
        st.write("")
        st.write("")
        if st.button("Approximate prices", key="btn_price"):
            if ballpark <= 0:
                st.warning("Enter a ballpark total above $0 first.")
            elif not P["milestones"]:
                st.warning("Add at least one milestone row first.")
            else:
                method = "duration" if method_label.startswith("Duration") else "even"
                distribute_ballpark(P["milestones"], ballpark, method=method)
                st.session_state["params"] = P
                st.session_state["price_done"] = True
                _bump_ms_editor()
                st.rerun()

    # ---- auto-computed total ----
    subtotal, n_parsed, n_unparsed = compute_total(P["milestones"])
    tc1, tc2, tc3 = st.columns([1, 1, 2])
    with tc1:
        st.metric("Project total (auto)", format_money(subtotal))
    with tc2:
        add_star = st.checkbox("Append \u201c*\u201d (footnote ref)",
                               value=P.get("total_asterisk", True))
        P["total_asterisk"] = add_star
    with tc3:
        P["total_note"] = st.text_input("Slide 7 footnote", P.get("total_note", ""))
    P["total_price"] = format_money(subtotal) + ("*" if add_star else "")
    if n_unparsed:
        st.caption(f"\u26a0\ufe0f {n_unparsed} price cell(s) couldn't be read as a "
                   f"number and were excluded from the total. Use values like "
                   f"`2500`, `$2,500`, or `1400`.")

    # ====================================================================== #
    # STEP 4 — Generate the deck
    # ====================================================================== #
    st.subheader("4) Generate proposal deck")
    _on_deck = (mode == "Revise existing") and bool(st.session_state.get("base_deck_bytes"))
    if _on_deck:
        st.caption("Build mode: **revising your uploaded deck** \u2014 manual formatting is "
                   "preserved; only changed content is rewritten. The result can be "
                   "re-uploaded above for another round.")
    if st.button("Build .pptx", type="primary"):
        try:
            data = None
            if _on_deck:
                data = revise_proposal_pptx(
                    P, st.session_state["base_deck_bytes"],
                    prev_params=st.session_state.get("base_deck_params"))
                # the freshly built revision becomes the base for the next round
                st.session_state["base_deck_bytes"] = data
                st.session_state["base_deck_params"] = copy.deepcopy(P)
            elif not os.path.exists(TEMPLATE_PATH):
                st.error(f"Template '{TEMPLATE_FILENAME}' not found next to the app.")
            else:
                data = build_proposal_pptx(P, TEMPLATE_PATH)
            if data:
                safe = "".join(c for c in P["customer_short"] if c.isalnum() or c in " _-").strip().replace(" ", "_")
                rev_tag = ""
                if _on_deck:
                    rv = (_read_state_from_bytes(data) or {}).get("_revision", 0)
                    rev_tag = f"_rev{rv}"
                fname = f"BasePair_Proposal_{safe}_{datetime.now().strftime('%Y%m%d_%H%M')}{rev_tag}.pptx"
                st.download_button("Download proposal deck", data=data, file_name=fname,
                                   mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
                st.success("Revision built onto your deck." if _on_deck else "Deck built.")
        except Exception as e:
            st.error(f"Build failed: {e}")

    st.session_state["params"] = P
    st.caption(f"{os.path.basename(__file__)} ({APP_VERSION}) "
               f"\u2014 {datetime.now().strftime('%Y-%m-%d %H:%M')}")


if __name__ == "__main__":
    run_app()
